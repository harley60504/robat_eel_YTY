from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parent / "output" / "eel_scan_rl_automation_flow_v15.pptx"
PREVIEW_IMAGE = Path(r"C:\Users\ytyla\Pictures\real_vs_mujoco_all.png")
SWIM_PREVIEWS = [
    ("直線游泳預覽", Path(r"C:\Users\ytyla\Pictures\straight_real_vs_mujoco.png")),
    ("左轉游泳預覽", Path(r"C:\Users\ytyla\Pictures\turn_left_real_vs_mujoco.png")),
    ("原地左旋預覽", Path(r"C:\Users\ytyla\Pictures\spin_left_real_vs_mujoco.png")),
]

W, H = Inches(13.333), Inches(7.5)

INK = RGBColor(28, 36, 43)
MUTED = RGBColor(92, 104, 112)
PAPER = RGBColor(248, 246, 241)
LINE = RGBColor(193, 184, 171)
BLUE = RGBColor(36, 103, 159)
TEAL = RGBColor(35, 145, 130)
GOLD = RGBColor(204, 139, 42)
RED = RGBColor(183, 76, 67)
WHITE = RGBColor(255, 255, 255)
SOFT = RGBColor(238, 234, 226)


def add_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = PAPER


def tx(slide, text, x, y, w, h, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft JhengHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def kicker(slide, label, no):
    tx(slide, f"{no:02d}  {label}", 0.55, 0.38, 2.4, 0.25, 9, BLUE, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.72), Inches(0.58), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def title(slide, claim, sub=None):
    tx(slide, claim, 0.55, 0.82, 11.7, 0.72, 28, INK, True)
    if sub:
        tx(slide, sub, 0.58, 1.55, 10.5, 0.36, 13, MUTED)


def footer(slide, no):
    tx(slide, "Robot eel MuJoCo | Hopf CPG parameter scan + RL fine tuning", 0.55, 7.12, 7.2, 0.2, 8, MUTED)
    tx(slide, str(no), 12.45, 7.12, 0.4, 0.2, 8, MUTED, align=PP_ALIGN.RIGHT)


def pill(slide, text, x, y, w, h, fill=WHITE, border=LINE, color=INK, size=14, bold=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(1.1)
    tf = shp.text_frame
    tf.clear()
    tf.margin_left = Inches(0.13)
    tf.margin_right = Inches(0.13)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft JhengHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return shp


def rect(slide, text, x, y, w, h, fill=WHITE, border=LINE, color=INK, size=13, bold=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(1)
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.13)
    tf.margin_right = Inches(0.13)
    tf.margin_top = Inches(0.09)
    tf.margin_bottom = Inches(0.09)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft JhengHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return shp


def arrow(slide, x1, y1, x2, y2, color=BLUE, width=2.0):
    con = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    con.line.color.rgb = color
    con.line.width = Pt(width)
    con.line.end_arrowhead = True
    return con


def line(slide, x1, y1, x2, y2, color=LINE, width=1.2):
    con = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    con.line.color.rgb = color
    con.line.width = Pt(width)
    return con


def add_picture_contained(slide, image_path, x, y, max_w, max_h):
    from PIL import Image

    with Image.open(image_path) as im:
        img_w, img_h = im.size
    scale = min(max_w / img_w, max_h / img_h)
    w = img_w * scale
    h = img_h * scale
    px = x + (max_w - w) / 2
    py = y + (max_h - h) / 2
    return slide.shapes.add_picture(str(image_path), Inches(px), Inches(py), width=Inches(w), height=Inches(h))


def bullet(slide, items, x, y, w, h, size=14, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(7)
        r = p.add_run()
        r.text = f"• {item}"
        r.font.name = "Microsoft JhengHei"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def mini_wave(slide, x, y, w, h):
    pts = []
    for i in range(6):
        cx = x + w * i / 5
        cy = y + h * (0.5 + (0.28 if i % 2 else -0.28))
        pts.append((cx, cy))
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.055), Inches(cy - 0.055), Inches(0.11), Inches(0.11))
        circ.fill.solid()
        circ.fill.fore_color.rgb = TEAL
        circ.line.fill.background()
    for a, b in zip(pts, pts[1:]):
        line(slide, a[0], a[1], b[0], b[1], TEAL, 2)
    tx(slide, "6 joints: amp scales + phase lags", x, y + h + 0.1, w, 0.2, 9, MUTED, align=PP_ALIGN.CENTER)


def rectangle_course(slide, x, y, w, h):
    line(slide, x, y, x + w, y, BLUE, 2.2)
    line(slide, x + w, y, x + w, y + h, BLUE, 2.2)
    line(slide, x + w, y + h, x, y + h, BLUE, 2.2)
    line(slide, x, y + h, x, y, BLUE, 2.2)
    pts = [(x + w, y), (x + w, y + h), (x, y + h), (x, y)]
    labels = ["P1", "P2", "P3", "P4"]
    for (px, py), label in zip(pts, labels):
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px - 0.08), Inches(py - 0.08), Inches(0.16), Inches(0.16))
        c.fill.solid()
        c.fill.fore_color.rgb = GOLD
        c.line.fill.background()
        tx(slide, label, px - 0.22, py - 0.34, 0.45, 0.18, 8, MUTED, True, PP_ALIGN.CENTER)
    fish_x = x + w * 0.48
    fish_y = y + h * 0.82
    body = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(fish_x - 0.09), Inches(fish_y - 0.07), Inches(0.18), Inches(0.14))
    body.rotation = 270
    body.fill.solid()
    body.fill.fore_color.rgb = TEAL
    body.line.fill.background()
    tail = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(fish_x + 0.06), Inches(fish_y - 0.06), Inches(0.12), Inches(0.12))
    tail.rotation = 90
    tail.fill.solid()
    tail.fill.fore_color.rgb = TEAL
    tail.line.fill.background()
    tx(slide, "robot eel", fish_x - 0.34, fish_y + 0.13, 0.7, 0.18, 7, TEAL, True, PP_ALIGN.CENTER)
    tx(slide, "3 m x 1.5 m rectangle course", x - 0.15, y + h + 0.23, w + 0.3, 0.24, 10, MUTED, align=PP_ALIGN.CENTER)


def add_cpg_algorithm_slide(prs, blank, no):
    s = prs.slides.add_slide(blank)
    add_bg(s); kicker(s, "CPG 演算法", no); title(s, "Hopf CPG 先定義機器魚怎麼擺尾，後面掃描與 RL 都是在調它。", "每個關節是一個 Hopf oscillator；相鄰關節透過相位耦合形成往後傳的波，最後輸出 servo 目標角。")
    rect(s, "輸入參數\nf, λ, ajoint\nmu_scales, phase_lags\njoint_bias", 0.85, 2.0, 2.15, 1.25, WHITE, BLUE, BLUE, 12, True)
    arrow(s, 3.15, 2.63, 4.0, 2.63, BLUE, 1.8)
    rect(s, "相位偏移\noffset[0] = 0\noffset[j] = -Σ phase_lags", 4.1, 2.0, 2.45, 1.25, WHITE, TEAL, TEAL, 12, True)
    arrow(s, 6.7, 2.63, 7.55, 2.63, TEAL, 1.8)
    rect(s, "狀態更新\nr[j], θ[j]\nEuler integration", 7.65, 2.0, 2.2, 1.25, WHITE, GOLD, GOLD, 12, True)
    arrow(s, 9.98, 2.63, 10.83, 2.63, GOLD, 1.8)
    rect(s, "輸出角度\nq[j] = ajoint r[j] cos θ[j] + bias[j]", 10.93, 2.0, 2.2, 1.25, WHITE, RED, RED, 11, True)
    tx(s, "核心更新式", 0.9, 3.85, 2.0, 0.3, 16, INK, True)
    rect(s, "ω = 2πf\nμ_target[j] = μ × mu_scale[j]\ndr[j] = α(μ_target[j] - r[j]^2)r[j]", 0.95, 4.35, 3.7, 1.25, INK, INK, WHITE, 12, False)
    rect(s, "dθ[j] = ω\n+ k_couple sin(phase error to neighbors)\n+ k_anchor sin(reference phase error)", 4.95, 4.35, 3.85, 1.25, INK, INK, WHITE, 12, False)
    rect(s, "r ← max(0, r + dr × dt)\nθ ← wrap_pi(θ + dθ × dt)\nservo target ← clip(q, -1.2, 1.2)", 9.1, 4.35, 3.2, 1.25, INK, INK, WHITE, 12, False)
    bullet(s, ["掃描調的是 `mu_scales` 與 `phase_lags`，等於調整波形的振幅分布與節間相位。", "矩形路徑控制額外改 `joint_bias` 與轉彎時的 `mu_scales`，不需要重寫 CPG。"], 0.95, 6.02, 11.4, 0.72, 12)
    footer(s, no)


def add_rl_tuning_slide(prs, blank, no):
    s = prs.slides.add_slide(blank)
    add_bg(s); kicker(s, "RL 調哪裡", no); title(s, "RL 調的是 CPG 的形狀參數，不是直接輸出每顆馬達角度。", "在目前 `EelFreeSwimRLEnv`，action 會被轉成 6 個振幅倍率與 5 個相位差，再塞回 HopfCPGParams。")
    rect(s, "RL action\n11 維 normalized action\n[-1, 1]", 0.8, 2.15, 2.25, 0.95, WHITE, BLUE, BLUE, 13, True)
    arrow(s, 3.18, 2.63, 4.05, 2.63, BLUE, 1.8)
    rect(s, "_physical_action()\n把 [-1,1] 映射到 bounds", 4.15, 2.15, 2.45, 0.95, WHITE, TEAL, TEAL, 13, True)
    arrow(s, 6.72, 2.63, 7.58, 2.63, TEAL, 1.8)
    rect(s, "amp_scales[1..6]\n控制每節擺尾振幅", 7.68, 2.15, 2.25, 0.95, WHITE, GOLD, GOLD, 13, True)
    arrow(s, 10.05, 2.63, 10.88, 2.63, GOLD, 1.8)
    rect(s, "phase_lags[1..5]\n控制波往後傳的相位差", 10.98, 2.15, 2.05, 0.95, WHITE, RED, RED, 12, True)
    tx(s, "塞回 CPG 的位置", 0.9, 3.82, 2.3, 0.3, 16, INK, True)
    rect(s, "HopfCPGParams(\n  frequency = fixed_frequency,\n  wavelength = fixed_wavelength,\n  ajoint = fixed_ajoint,\n  mu_scales = amp_scales_to_mu_scales(amp_scales),\n  phase_lags = phase_lags,\n)", 0.95, 4.28, 5.3, 1.75, INK, INK, WHITE, 12, False)
    rect(s, "所以 RL 實際學到的是：\n1. 哪幾節要擺大或擺小\n2. 每節波峰要慢一點或快一點出現\n3. 在直游目標下，如何同時壓漂移、yaw、能耗和平滑度", 6.75, 4.28, 5.45, 1.75, WHITE, BLUE, INK, 13, True)
    bullet(s, ["沒有 RL：掃描只會找到固定的一組參數。", "有 RL：policy 可根據 observation，在每一步重新給一組 CPG shape 參數。", "但它仍受掃描後 bounds 限制，所以比較不容易跑到奇怪的 gait。"], 0.95, 6.32, 11.6, 0.72, 11)
    footer(s, no)


def add_reward_to_cpg_detail_slide(prs, blank, no):
    s = prs.slides.add_slide(blank)
    add_bg(s); kicker(s, "Reward 到 CPG", no); title(s, "Reward 改 CPG 的重點鏈路：先改 policy theta，再由 action 寫入 CPG 參數。", "reward 算分 -> 更新 policy theta -> policy 輸出 action -> _physical_action() 轉 bounds -> amp_scales / phase_lags -> HopfCPGParams")

    steps = [
        ("1", "reward 算分", "MuJoCo 跑一步後量測 vx、漂移、yaw、能耗、action 變化。\n分數高代表這個 action 讓 eel 更直、更穩、更省。", BLUE),
        ("2", "更新 policy theta", "RL 演算法用 reward 調整神經網路權重 theta。\n高 reward 的 action 機率提高，低 reward 的 action 機率降低。", TEAL),
        ("3", "policy 輸出 action", "新的 theta 讀 observation 後輸出 11 維 action。\naction 仍是 normalized [-1, 1]，不是實際馬達角度。", GOLD),
    ]
    for i, (num, head, body, col) in enumerate(steps):
        x = 0.72 + i * 4.12
        pill(s, num, x, 2.02, 0.42, 0.42, col, col, WHITE, 13, True)
        tx(s, head, x + 0.52, 2.04, 2.8, 0.28, 15, col, True)
        rect(s, body, x, 2.55, 3.35, 1.05, WHITE, col, INK, 10, False)
        if i < 2:
            arrow(s, x + 3.42, 3.08, x + 3.92, 3.08, col, 1.6)

    tx(s, "執行時如何真正改到 CPG", 0.82, 4.05, 3.5, 0.3, 16, INK, True)
    rect(s, "action[0:6]\n-> amp_scales\n-> amp_scales_to_mu_scales()\n-> HopfCPGParams.mu_scales\n\n效果：改每一節的振幅大小；某節擺大、某節擺小，身體波形就改變。", 0.82, 4.48, 3.65, 1.55, INK, INK, WHITE, 10, False)
    rect(s, "action[6:11]\n-> phase_lags\n-> HopfCPGParams.phase_lags\n\n效果：改相鄰關節波峰出現的時間差；相位差變大/變小，推進波往尾巴傳的節奏就改變。", 4.82, 4.48, 3.65, 1.55, INK, INK, WHITE, 10, False)
    rect(s, "HopfCPGParams(\n  fixed frequency / wavelength / ajoint,\n  mu_scales = new_mu_scales,\n  phase_lags = new_phase_lags,\n)\n\n下一個控制步：Hopf CPG 用新參數積分，產生新的 q[j] servo target。", 8.82, 4.48, 3.55, 1.55, WHITE, RED, INK, 10, True)
    arrow(s, 4.48, 5.25, 4.78, 5.25, BLUE, 1.4)
    arrow(s, 8.48, 5.25, 8.78, 5.25, GOLD, 1.4)
    tx(s, "所以 reward 的作用不是「直接改 CPG 公式」，而是改 policy theta；policy theta 改變後，下一次 action 才會把新的振幅與相位寫進 CPG。", 0.9, 6.35, 11.35, 0.42, 13, RED, True, PP_ALIGN.CENTER)
    footer(s, no)


def add_ppo_rationale_slide(prs, blank, no):
    s = prs.slides.add_slide(blank)
    add_bg(s); kicker(s, "Why PPO", no); title(s, "不用 Q-learning 查表，因為 CPG action 是 11 維連續控制。", "PPO 直接學 policy theta，讓 observation 對應到連續的 amp_scales 與 phase_lags。")

    tx(s, "Q-learning 的問題", 0.7, 1.92, 2.6, 0.28, 15, RED, True)
    rect(s, "Q-learning / DQN 核心是學 Q(s, a)：\n在某 state 做某 action，未來總分大概多少。\n\n但本研究 action 是：\n6 個 amp_scales + 5 個 phase_lags\n= 11 維連續 action。\n\n若每一維只切 5 格：\n5^11 = 48,828,125 種組合。\n切 10 格則是 10^11 種。\n\n結果：查表或離散化會讓 action space 爆炸，而且會破壞 CPG 參數的平滑微調。", 0.72, 2.32, 3.55, 3.62, WHITE, RED, INK, 10, False)

    tx(s, "PPO 的做法", 4.55, 1.92, 2.4, 0.28, 15, BLUE, True)
    rect(s, "PPO 是 policy-gradient / actor-critic 方法。\n它不需要列出所有 action，而是直接學：\n\npolicy theta(observation) -> action\n\n在本系統中：\nobservation -> 11 維 action\n-> _physical_action()\n-> amp_scales / phase_lags\n-> HopfCPGParams\n\nPPO 用 clipped objective 限制 policy 更新幅度，讓訓練比一般 policy gradient 穩定，也比 TRPO 更容易實作與調參。", 4.55, 2.32, 3.55, 3.62, WHITE, BLUE, INK, 10, False)

    rect(s, "放在本研究的解釋：\n\nQ-learning 像是在 state-action 空間裡替每個候選動作估分；\nPPO 則是直接調整 policy theta，讓 policy 輸出更好的連續 CPG action。\n\n因此 reward 的作用路徑是：\nreward / advantage -> policy theta -> continuous action -> CPG shape。", 8.45, 2.32, 3.65, 3.62, SOFT, LINE, INK, 11, False)

    tx(s, "論文依據：Schulman et al., Proximal Policy Optimization Algorithms, arXiv:1707.06347, 2017.  對照：Mnih et al., Human-level control through deep reinforcement learning, Nature, 2015.", 0.76, 6.31, 11.7, 0.22, 8, MUTED)
    tx(s, "CPG + DRL 仿生游泳根據：Hameed et al., Deep reinforcement learning enabling a BCFbot to learn various undulatory patterns, Ocean Engineering 320, 120322, 2025. DOI: 10.1016/j.oceaneng.2025.120322.", 0.76, 6.55, 11.7, 0.22, 8, MUTED)
    footer(s, no)


def add_swim_preview_slide(prs, blank, no, label, image_path):
    s = prs.slides.add_slide(blank)
    add_bg(s); kicker(s, "實際游泳預覽", no); title(s, f"{label}：用實際游泳與 MuJoCo 對照確認 gait 行為。")
    if image_path.exists():
        add_picture_contained(s, image_path, 1.15, 1.75, 10.95, 4.82)
        tx(s, f"Source: {image_path}", 0.95, 6.78, 11.2, 0.22, 9, MUTED, align=PP_ALIGN.CENTER)
    else:
        rect(s, f"Preview image missing:\n{image_path}", 1.4, 2.3, 10.5, 2.4, SOFT, LINE, MUTED, 16, True)
    footer(s, no)


def add_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # 1 cover
    s = prs.slides.add_slide(blank)
    add_bg(s)
    tx(s, "Robot Eel", 0.55, 0.55, 2.8, 0.35, 15, BLUE, True)
    tx(s, "先掃描直線參數，再用 RL 微修", 0.55, 1.18, 9.4, 0.85, 34, INK, True)
    tx(s, "用粗搜尋先找到穩定直游的 Hopf CPG 參數區間，再讓強化學習在小範圍內修正漂移、yaw、能耗與平滑度。", 0.6, 2.16, 8.4, 0.72, 16, MUTED)
    mini_wave(s, 7.3, 3.15, 4.1, 1.2)
    rect(s, "掃描輸出\nbest params / top rows", 0.9, 4.35, 2.25, 0.95, WHITE, BLUE, BLUE, 13, True)
    arrow(s, 3.2, 4.83, 4.25, 4.83)
    rect(s, "縮小 RL action bounds\namp + phase window", 4.35, 4.35, 2.55, 0.95, WHITE, TEAL, TEAL, 13, True)
    arrow(s, 6.95, 4.83, 7.95, 4.83, TEAL)
    rect(s, "RL 微調\nreward-weighted policy", 8.05, 4.35, 2.45, 0.95, WHITE, GOLD, GOLD, 13, True)
    arrow(s, 10.55, 4.83, 11.38, 4.83, GOLD)
    rect(s, "驗證與部署\nstraight swim gait", 11.45, 4.35, 1.55, 0.95, WHITE, RED, RED, 11, True)
    footer(s, 1)

    # 2 CPG algorithm first
    add_cpg_algorithm_slide(prs, blank, 2)

    # 3 core idea
    s = prs.slides.add_slide(blank)
    add_bg(s); kicker(s, "核心想法", 3); title(s, "不要讓 RL 從整片海裡找答案，先把可用參數區間掃出來。", "掃描負責找到可前進且不太偏的初始 gait；RL 負責在可行區內做細節修正。")
    rect(s, "直接 RL", 0.8, 2.35, 2.1, 0.7, SOFT, LINE, INK, 16, True)
    arrow(s, 3.0, 2.7, 4.55, 2.7, RED)
    rect(s, "探索空間大\n容易學到不穩定或偏航策略", 4.7, 2.12, 3.0, 1.15, WHITE, RED, RED, 14, True)
    rect(s, "掃描 + RL", 0.8, 4.25, 2.1, 0.7, SOFT, LINE, INK, 16, True)
    arrow(s, 3.0, 4.6, 4.55, 4.6, BLUE)
    rect(s, "先取得穩定直線候選\n再微調性能與魯棒性", 4.7, 4.02, 3.0, 1.15, WHITE, BLUE, BLUE, 14, True)
    bullet(s, ["掃描：可平行、可重現、容易知道哪組參數好。", "RL：只在掃描後的窄範圍內微調，樣本效率比較高。", "最後再用同一套量測腳本驗證直線速度、側向漂移與 yaw。"], 8.25, 2.15, 4.2, 3.3, 15)
    footer(s, 3)

    # 4 parameter scan
    s = prs.slides.add_slide(blank)
    add_bg(s); kicker(s, "參數掃描", 4); title(s, "掃描階段搜尋的是 CPG 形狀，而不是直接控制馬達。", "目前程式把 Hopf CPG 的輸出轉成 6 個尾部關節目標角，掃描的核心是振幅比例與相位差。")
    tx(s, "被掃描的 11 維參數", 0.8, 2.15, 3.6, 0.35, 16, INK, True)
    for i in range(6):
        rect(s, f"A{i+1}", 0.85 + i*0.68, 2.85, 0.52, 0.48, WHITE, TEAL, TEAL, 14, True)
    for i in range(5):
        rect(s, f"φ{i+1}", 1.2 + i*0.82, 3.72, 0.58, 0.48, WHITE, BLUE, BLUE, 14, True)
    mini_wave(s, 6.8, 2.35, 4.2, 1.0)
    arrow(s, 4.98, 3.28, 6.35, 3.28, BLUE)
    rect(s, "HopfCPGParams\nfrequency / wavelength / ajoint\nmu_scales / phase_lags", 6.25, 4.1, 3.1, 1.25, WHITE, LINE, INK, 13, True)
    arrow(s, 9.45, 4.72, 10.5, 4.72, TEAL)
    rect(s, "MuJoCo simulation\nmeasure Fx, Fy, energy", 10.6, 4.1, 2.25, 1.25, WHITE, LINE, INK, 13, True)
    bullet(s, ["`amp_scales_to_mu_scales()` 把振幅倍率平方後餵給 Hopf 振盪器。", "`phase_lags` 決定每節之間的波往後傳遞多少。", "掃描不是盲目微調 servo，而是找一組可解釋的游動波形。"], 0.85, 4.75, 5.4, 1.4, 13)
    footer(s, 4)

    # 5 scoring
    s = prs.slides.add_slide(blank)
    add_bg(s); kicker(s, "評分邏輯", 5); title(s, "直線 gait 的第一關，是把前向推力做大，同時壓住側向與能耗。")
    rect(s, "score = mean_fx - fy_weight × |mean_fy| - energy_weight × energy_proxy", 1.0, 2.0, 11.1, 0.75, INK, INK, WHITE, 17, True)
    labels = [("mean_fx", "想要大", BLUE, 0.80), ("|mean_fy|", "想要小", RED, 0.35), ("energy", "想要小", GOLD, 0.42)]
    for idx, (name, cap, col, val) in enumerate(labels):
        x = 1.15 + idx * 3.85
        tx(s, name, x, 3.25, 2.8, 0.3, 18, col, True, PP_ALIGN.CENTER)
        h = 1.65 * val
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.95), Inches(5.45 - h), Inches(0.65), Inches(h))
        r.fill.solid(); r.fill.fore_color.rgb = col; r.line.fill.background()
        tx(s, cap, x, 5.65, 2.8, 0.25, 12, MUTED, align=PP_ALIGN.CENTER)
    bullet(s, ["保留 top rows：不是只存一個 best，而是留一批候選給後續比較。", "verify-seconds：最佳參數需要跑更久，避免只是在短時間偶然高分。", "輸出 CSV/JSON 後，下一步可自動轉成 RL bounds。"], 1.0, 6.05, 11.2, 0.78, 12)
    footer(s, 5)

    # 6 RL
    s = prs.slides.add_slide(blank)
    add_bg(s); kicker(s, "RL 微修", 6); title(s, "RL 不再負責找大方向，而是微修掃描後的可行 gait。", "EelFreeSwimRLEnv 的 action 仍是 6 個振幅倍率 + 5 個相位差，reward 多加了漂移、yaw、能耗和平滑度。")
    rect(s, "掃描 top-k\nbest rows", 0.9, 2.65, 2.0, 0.85, WHITE, BLUE, BLUE, 14, True)
    arrow(s, 3.02, 3.07, 4.15, 3.07, BLUE)
    rect(s, "轉成 action bounds\nlows / highs", 4.25, 2.65, 2.3, 0.85, WHITE, TEAL, TEAL, 14, True)
    arrow(s, 6.68, 3.07, 7.78, 3.07, TEAL)
    rect(s, "RL policy learns\nfine adjustment", 7.88, 2.65, 2.45, 0.85, WHITE, GOLD, GOLD, 14, True)
    arrow(s, 10.45, 3.07, 11.28, 3.07, GOLD)
    rect(s, "驗證直游\nspeed / drift / yaw", 11.35, 2.65, 1.65, 0.85, WHITE, RED, RED, 12, True)
    tx(s, "reward terms", 0.9, 4.45, 2.0, 0.3, 16, INK, True)
    for i, (t, c) in enumerate([("forward vx", BLUE), ("lateral vy / y", RED), ("yaw / yaw rate", GOLD), ("energy", TEAL), ("smooth action", MUTED)]):
        pill(s, t, 0.95 + i*2.43, 5.0, 2.05, 0.52, WHITE, c, c, 12, True)
    bullet(s, ["掃描提供合理初始區域：RL action 不需要涵蓋所有可能。", "reward 讓目標從「推得動」變成「游得直、穩、低能耗」。", "若換目標，例如矩形路徑或轉彎，只要替換 reward / bounds 來源即可。"], 0.95, 6.0, 11.6, 0.85, 12)
    footer(s, 6)

    # 7 where RL tunes CPG
    add_rl_tuning_slide(prs, blank, 7)

    # 8 how reward changes CPG through policy
    add_reward_to_cpg_detail_slide(prs, blank, 8)

    # 9 why PPO instead of Q-learning
    add_ppo_rationale_slide(prs, blank, 9)

    # 10 automation flowchart
    s = prs.slides.add_slide(blank)
    add_bg(s); kicker(s, "自動化流程", 10); title(s, "整個流程可以做成一次執行：掃描、收斂、RL、驗證、輸出報告。")
    nodes = [
        ("設定基準參數\nfreq / wavelength / ajoint", 0.75, 2.05, BLUE),
        ("產生候選\nrandom starts + local steps", 3.05, 2.05, TEAL),
        ("MuJoCo 評估\nFx / Fy / energy", 5.35, 2.05, GOLD),
        ("排序與保存\ntop-k CSV / JSON", 7.65, 2.05, BLUE),
        ("轉 RL bounds\nnarrow action space", 9.95, 2.05, TEAL),
        ("訓練 RL\nreward shaping", 2.0, 4.65, GOLD),
        ("驗證直線游\nspeed / drift / yaw", 4.9, 4.65, RED),
        ("輸出最佳 gait\nparams + plots + log", 7.8, 4.65, BLUE),
    ]
    for text, x, y, col in nodes:
        rect(s, text, x, y, 1.9, 0.9, WHITE, col, col, 11, True)
    for x1, x2 in [(2.65, 3.05), (4.95, 5.35), (7.25, 7.65), (9.55, 9.95)]:
        arrow(s, x1, 2.5, x2, 2.5, BLUE, 1.8)
    arrow(s, 10.9, 2.95, 2.95, 4.65, TEAL, 1.8)
    arrow(s, 3.9, 5.1, 4.9, 5.1, GOLD, 1.8)
    arrow(s, 6.8, 5.1, 7.8, 5.1, RED, 1.8)
    arrow(s, 8.55, 4.65, 8.1, 2.95, BLUE, 1.5)
    tx(s, "未達門檻就回到排序/範圍設定", 8.45, 3.62, 2.5, 0.3, 10, MUTED, align=PP_ALIGN.CENTER)
    footer(s, 10)

    # 11 implementation outputs
    s = prs.slides.add_slide(blank)
    add_bg(s); kicker(s, "實作輸出", 11); title(s, "自動化腳本應該留下三種東西：參數、軌跡、決策依據。")
    rect(s, "1. 參數檔\nbest_params.json\namp_scales / phase_lags", 0.9, 2.2, 3.0, 1.2, WHITE, BLUE, BLUE, 14, True)
    rect(s, "2. 指標表\nscan_results.csv\nscore / Fx / Fy / energy", 4.25, 2.2, 3.0, 1.2, WHITE, TEAL, TEAL, 14, True)
    rect(s, "3. 驗證軌跡\ntrajectory.csv / plots\nx-y / yaw / velocity", 7.6, 2.2, 3.0, 1.2, WHITE, GOLD, GOLD, 14, True)
    bullet(s, ["掃描後：把 top-k 的 min/max 或 best ± margin 轉為 RL bounds。", "RL 後：用相同 measurement script 重跑，避免 reward 和真實評估脫節。", "最終報告：列出 baseline、scan-best、RL-best 三者比較。"], 0.95, 4.2, 5.9, 1.7, 14)
    rect(s, "建議指令範例\npython auto_tune_hopf_shape.py --xml eel_tethered.xml --output-dir outputs/hopf_shape_auto\npython measure_free_swim_speed.py --amp-scales ... --phase-lags ... --csv outputs/verify.csv", 7.2, 4.05, 5.2, 1.95, INK, INK, WHITE, 11, False)
    footer(s, 11)

    # 12 rectangle path following
    s = prs.slides.add_slide(blank)
    add_bg(s); kicker(s, "繞長方形流程", 12); title(s, "直線 gait 找好後，長方形路徑只需要在上層加 steering。", "底層仍用掃描/RL 得到的 Hopf CPG 直游參數；上層 pure pursuit 決定目前要往哪個方向修正。")
    rectangle_course(s, 10.35, 2.1, 1.95, 1.18)
    rect(s, "1. 取得位置\nbase x, y, yaw", 0.8, 2.22, 2.1, 0.74, WHITE, BLUE, BLUE, 12, True)
    arrow(s, 2.98, 2.59, 3.52, 2.59, BLUE, 1.8)
    rect(s, "2. 最近路徑進度\nclosest_s(x, y)", 3.62, 2.22, 2.15, 0.74, WHITE, TEAL, TEAL, 12, True)
    arrow(s, 5.86, 2.59, 6.42, 2.59, TEAL, 1.8)
    rect(s, "3. 前視目標點\npoint_at(s + lookahead)", 6.52, 2.22, 2.35, 0.74, WHITE, GOLD, GOLD, 11, True)
    arrow(s, 7.7, 3.02, 2.08, 4.12, GOLD, 1.4)
    rect(s, "4. heading error\nwrap_pi(desired_yaw - yaw)", 0.8, 4.2, 2.55, 0.78, WHITE, RED, RED, 11, True)
    arrow(s, 3.45, 4.59, 4.02, 4.59, RED, 1.8)
    rect(s, "5. steering bias\nlow-pass + clamp", 4.12, 4.2, 2.25, 0.78, WHITE, BLUE, BLUE, 12, True)
    arrow(s, 6.48, 4.59, 7.02, 4.59, BLUE, 1.8)
    rect(s, "6. 更新 CPG\njoint_bias + turn amp", 7.12, 4.2, 2.35, 0.78, WHITE, TEAL, TEAL, 12, True)
    arrow(s, 9.58, 4.59, 10.12, 4.59, TEAL, 1.8)
    rect(s, "7. 評估\nlaps / contacts / bounds", 10.22, 4.2, 2.1, 0.78, WHITE, GOLD, GOLD, 11, True)
    bullet(s, ["`pure_pursuit`：用最近路徑進度加上 `lookahead`，不是只追下一個角點。", "`steering_profile()`：把單一 steer 值變成 6 個 joint bias，尾端修正較大。", "`turning_amp_scales()`：轉彎時尾端振幅增加，讓身體更容易繞角。"], 0.95, 5.65, 11.8, 0.9, 12)
    footer(s, 12)

    # 13 conclusion
    s = prs.slides.add_slide(blank)
    add_bg(s); kicker(s, "結論", 13); title(s, "這個流程的價值，是把可解釋搜尋和資料驅動微調接起來。")
    for i, (head, body, col) in enumerate([
        ("先懂 CPG", "CPG 定義擺尾波形；掃描與 RL 都是在調這個波形。", BLUE),
        ("再掃描", "用物理模擬快速找出可直游的 CPG 參數區域。", TEAL),
        ("最後微修與繞路", "RL 修正直游品質，長方形路徑再用 steering 接上。", GOLD),
    ]):
        x = 1.05 + i * 4.05
        pill(s, head, x, 2.35, 2.05, 0.55, WHITE, col, col, 16, True)
        tx(s, body, x - 0.35, 3.15, 2.75, 0.9, 15, INK, align=PP_ALIGN.CENTER)
    tx(s, "重點：CPG 負責產生擺尾波形，掃描先找穩定直游參數，RL 再在合理範圍內微調，最後用 steering 接上長方形路徑。", 1.2, 5.15, 11.0, 0.8, 18, INK, True, PP_ALIGN.CENTER)
    footer(s, 13)

    # 14-16 real swimming previews
    for offset, (label, image_path) in enumerate(SWIM_PREVIEWS):
        add_swim_preview_slide(prs, blank, 14 + offset, label, image_path)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = add_deck()
    print(path)
